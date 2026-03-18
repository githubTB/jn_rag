from .base import BaseExtractor
from models.document import Document


class PptxExtractor(BaseExtractor):
    """Extract text from .pptx files, one Document per slide."""

    def __init__(self, file_path: str):
        self._file_path = file_path

    def extract(self) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(self._file_path)
        documents: list[Document] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)

            # Extract table content
            for shape in slide.shapes:
                if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    table = shape.table
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_texts:
                            texts.append(" | ".join(row_texts))

            if texts:
                documents.append(
                    Document(
                        page_content="\n".join(texts),
                        metadata={"source": self._file_path, "slide": slide_num},
                    )
                )

        return documents
